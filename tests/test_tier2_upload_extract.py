"""Tier 2 upload text extraction (slides, spreadsheets, CSV, RTF)."""

from __future__ import annotations

import tempfile
from pathlib import Path

from app.services.ingestion import ALLOWED_UPLOAD_EXTENSIONS, TIER2_UPLOAD_EXTENSIONS, extract_pages_from_upload


def test_tier2_extensions_union() -> None:
    assert ".pptx" in TIER2_UPLOAD_EXTENSIONS
    assert ".csv" in TIER2_UPLOAD_EXTENSIONS
    assert ".pdf" in ALLOWED_UPLOAD_EXTENSIONS
    assert ".pptx" in ALLOWED_UPLOAD_EXTENSIONS


def test_extract_csv() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "data.csv"
        p.write_text("col_a,col_b\nhello,world\n", encoding="utf-8")
        pages = extract_pages_from_upload(str(p), "data.csv")
        assert len(pages) == 1
        assert "hello" in pages[0].text and "world" in pages[0].text


def test_extract_rtf() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "note.rtf"
        p.write_text(r"{\rtf1\ansi Tier two RTF sample.\par}", encoding="utf-8")
        pages = extract_pages_from_upload(str(p), "note.rtf")
        assert len(pages) == 1
        assert "RTF" in pages[0].text or "Tier" in pages[0].text


def test_extract_xlsx_minimal() -> None:
    from openpyxl import Workbook

    with tempfile.TemporaryDirectory() as tmp:
        path = Path(tmp) / "book.xlsx"
        wb = Workbook()
        ws = wb.active
        assert ws is not None
        ws.title = "Demo"
        ws.append(["Key", "Value"])
        ws.append(["SVL", "Tier2"])
        wb.save(path)
        pages = extract_pages_from_upload(str(path), "book.xlsx")
        assert len(pages) >= 1
        joined = " ".join(p.text for p in pages)
        assert "SVL" in joined and "Tier2" in joined


def test_extract_pptx_minimal() -> None:
    from pptx import Presentation

    with tempfile.TemporaryDirectory() as tmp:
        path = Path(tmp) / "deck.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        assert tf is not None
        tf.text = "Sterling Vale slide body for tier two."
        prs.save(path)
        pages = extract_pages_from_upload(str(path), "deck.pptx")
        assert len(pages) >= 1
        assert "Sterling" in pages[0].text or "tier" in pages[0].text.lower()
