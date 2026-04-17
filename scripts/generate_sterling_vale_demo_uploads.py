"""
Generate Tier-1 demo files (PDF, DOCX, MD, HTML, TXT), Tier-2 samples (PPTX, XLSX, CSV, RTF,
optional XLS), plus `Sample_Questions.md` under
docs/demo/organizations/sterling-vale-llp/workspaces/<workspace-slug>/.

For manual upload testing against workspace-scoped RAG.

Usage (repo root):
  pip install -r requirements.txt
  pip install xlwt   # optional: adds Sterling_Vale_*_Tier2_Pack.xls

  python scripts/generate_sterling_vale_demo_uploads.py
"""

from __future__ import annotations

import csv
import html
import shutil
import subprocess
import sys
from email.message import EmailMessage
from pathlib import Path

_scripts = Path(__file__).resolve().parent
sys.path.insert(0, str(_scripts))

from fpdf import FPDF  # noqa: E402
from law_firm_demo_pdf_catalog import DEMO_ORGANIZATION_SLUG, workspace_dir_slug  # noqa: E402

_REPO = Path(__file__).resolve().parents[1]


def _fpdf_safe(text: str) -> str:
    """fpdf2 core fonts are Latin-1; normalize common Unicode punctuation."""
    for a, b in (
        ("\u2014", " - "),
        ("\u2013", "-"),
        ("\u00b7", "|"),
        ("\u2019", "'"),
    ):
        text = text.replace(a, b)
    return text

# Workspace key -> (display tag, long-form topic for filler)
_WORKSPACES: list[tuple[str, str, str]] = [
    ("General", "GENERAL", "firm-wide operations, retainers, and administrative coordination"),
    ("Litigation & Disputes", "LITIGATION", "motion practice, discovery, and docket control"),
    ("Corporate & Transactions", "CORPORATE", "vendor MSAs, indemnity, and deal closing"),
    ("Knowledge & Precedents", "KNOWLEDGE", "clause banks, template maintenance, and research memos"),
    ("Client Intake & Conflicts", "INTAKE", "new business intake, conflicts, and ethics walls"),
]

# Sample RAG / chat prompts grounded in generated Demo Pack themes (synthetic content only).
_SAMPLE_QUESTIONS_BY_TAG: dict[str, list[str]] = {
    "GENERAL": [
        "What file naming pattern does Sterling & Vale use instead of client names on shared drives?",
        "What three ingredients must align before work product circulates outside the firm according to the General demo pack?",
        "When should matter teams calendar internal checkpoints from file opening, and what is the example day pattern mentioned?",
        "Who should you contact when exposure estimates exceed the matter budget band?",
        "What prefix appears in the revision history table for demo ownership?",
        "What checklist item tells you to verify conflicts clearance timestamps, and what tag prefix is used?",
        "How should superseded drafts be archived according to the checklist?",
        "What does the demo text say about regulatory correspondence from an agency not on the intake form?",
    ],
    "LITIGATION": [
        "What topics does the Litigation demo pack focus on for motion practice and docket control?",
        "What escalation path is described when exposure exceeds the matter budget band?",
        "What document naming convention is repeated in the Litigation demo materials?",
        "According to the demo sections, what should matter teams calendar for quality review?",
        "What placeholder conflicts clearance format appears in the checklists (SVL-LITIGATION-####)?",
        "What should associates avoid putting in filenames on shared drives?",
        "What repeating section title pattern appears in the markdown demo (e.g. Practice note tied to the workspace topic)?",
        "Where should workspace retention labels appear before external circulation?",
    ],
    "CORPORATE": [
        "What deal-phase topics does the Corporate demo pack emphasize (vendor MSAs, indemnity, closing)?",
        "What SVL tag format is used in conflicts clearance verification lines in the demo pack?",
        "What internal reference pattern is repeated in segment lines across the Corporate files?",
        "When must partners align work product with the engagement letter and matter code?",
        "What filename pattern uses SVL-matter-doctype-version?",
        "What should teams do when a new regulator sends correspondence not listed on intake?",
        "What archive suffix is used for superseded drafts in the checklist?",
        "What does the demo text say about matter budget bands and escalation?",
    ],
    "KNOWLEDGE": [
        "What ongoing maintenance themes appear in the Knowledge & Precedents demo pack?",
        "What checklist references verifying conflicts clearance with SVL-KNOWLEDGE-####?",
        "What does the demo text say about clause banks and template maintenance?",
        "How should superseded language be tagged in the clause bank (hint: deprecated naming)?",
        "What revision-history owner role is listed in the demo markdown tables?",
        "What practice is described for research memos and internal reference segments?",
        "What document naming rule avoids client names on shared drives?",
        "What escalation is described when exposure exceeds the matter budget band?",
    ],
    "INTAKE": [
        "What intake themes are covered in the Client Intake demo pack (conflicts, ethics walls)?",
        "What timeline is mentioned for documenting a wall memo in the Client Intake workspace?",
        "What information should be collected about adverse parties in the intake narrative?",
        "What conflicts clearance format appears as SVL-INTAKE-#### in the demo checklists?",
        "When must work align with the engagement letter and matter code before external circulation?",
        "What should be avoided in filenames on shared drives per the demo text?",
        "Who should be contacted when exposure exceeds the matter budget band?",
        "What does the demo say about archived closed files in conflicts searches?",
    ],
}

# Tier 2 (slides, spreadsheets, CSV, RTF) — prompts reference Tier2_Pack filenames and embedded markers.
_SAMPLE_QUESTIONS_TIER2_BY_TAG: dict[str, list[str]] = {
    "GENERAL": [
        "In the Tier 2 PowerPoint, what appears in the title slide text starting with SVL Tier 2?",
        "What three worksheet tab names appear in the Tier 2 Excel workbook?",
        "What column header appears after workspace_tag in the Tier 2 CSV?",
        "What exact key string is written in cell row 3 column A of the Summary sheet (SVL-TIER2-GENERAL-...)?",
        "What phrase opens the Tier 2 RTF export body?",
        "What checkpoint day numbers appear in the body text on slide deck block 3?",
    ],
    "LITIGATION": [
        "What is the subtitle on slide 1 of the Tier 2 deck (tag and topic snippet)?",
        "In the Checklists sheet, what SVL-TIER2-LITIGATION marker appears in the key column?",
        "What tier2_marker value appears on CSV row for record SVL-LITIGATION-T2-0001?",
        "What does the Escalations sheet say about practice group leader contact?",
        "Does the RTF mention checkpoints or matter budget band?",
        "What is the title on slide deck block 2 in the PPTX?",
    ],
    "CORPORATE": [
        "What workspace name is in cell B2 of the Summary sheet in the Tier 2 xlsx?",
        "How many data rows (approx) are in the Tier 2 CSV after the header?",
        "What RTF text references vendor MSAs or indemnity themes from the topic line?",
        "What string identifies the Tier 2 XLS key row if you generated the .xls file?",
        "What appears in the title placeholder on the second content slide of the deck?",
    ],
    "KNOWLEDGE": [
        "Which sheet contains the phrase 'clause bank' or template maintenance in the Tier 2 spreadsheet?",
        "What is the first CSV record_id in the Tier 2 file?",
        "What slide title includes the KNOWLEDGE tag in the Tier 2 PowerPoint?",
        "What keyword appears in the RTF alongside workspace tag KNOWLEDGE?",
    ],
    "INTAKE": [
        "What timeline keyword appears in the Escalations sheet text (e.g. business days / wall memo)?",
        "What adversary-related theme from the topic appears in RTF or slide body text?",
        "What SVL-INTAKE CSV record_id format is used for row 5?",
        "What does slide deck block 1 title say in the Tier 2 PPTX?",
    ],
}

# Tier 3 (email + ebooks + OCR/image) prompts.
_SAMPLE_QUESTIONS_TIER3_BY_TAG: dict[str, list[str]] = {
    "GENERAL": [
        "What is the subject line in the Tier 3 EML sample for this workspace?",
        "What attachment marker appears in the EML body text (Attachment: ...)?",
        "What heading/title is extracted from the Tier 3 EPUB intro chapter?",
        "In the scanned PDF image sample, what SVL OCR marker string is visible?",
        "What workspace tag appears in the Tier 3 OCR image text block?",
    ],
    "LITIGATION": [
        "What does the Tier 3 EML mention about motion practice or docket control?",
        "What SVL OCR marker appears in the litigation scanned image PDF?",
        "What chapter heading is extracted from the Tier 3 EPUB file?",
        "What attachment filename is referenced in the email extraction output?",
    ],
    "CORPORATE": [
        "What vendor/MSA phrase appears in the Tier 3 EML body?",
        "Which SVL OCR marker appears in the corporate Tier 3 image sample?",
        "What line in the EPUB references indemnity/deal closing themes?",
        "What is the attachment filename included in the EML test pack?",
    ],
    "KNOWLEDGE": [
        "What clause bank / template maintenance phrase appears in Tier 3 EML?",
        "What OCR marker appears in the Knowledge scanned image sample?",
        "What chapter title appears in the Tier 3 EPUB content?",
        "What workspace tag appears in Tier 3 attachment text?",
    ],
    "INTAKE": [
        "What intake/conflicts phrase is present in the Tier 3 EML body?",
        "What OCR marker string appears in the intake scanned image sample?",
        "What timeline or wall memo phrase appears in Tier 3 EPUB text?",
        "What attachment marker is included in EML extraction output?",
    ],
}


def _sample_questions_markdown(
    workspace_name: str,
    tag: str,
    topic: str,
    stem: str,
    tier2_stem: str,
    tier3_stem: str,
    *,
    tier2_xls_written: bool,
    tier3_msg_written: bool,
    tier3_mobi_written: bool,
) -> str:
    questions = _SAMPLE_QUESTIONS_BY_TAG.get(tag, [])
    tier2_q = _SAMPLE_QUESTIONS_TIER2_BY_TAG.get(tag, [])
    tier3_q = _SAMPLE_QUESTIONS_TIER3_BY_TAG.get(tag, [])
    lines = [
        f"# Sample questions — {workspace_name}",
        "",
        f"**Workspace tag:** `{tag}` · **Demo focus:** {topic}",
        "",
        "Use these prompts **after** uploading demo files from this folder into the **same workspace** in the app. "
        "Synthetic internal reference text only — not legal advice.",
        "",
        "## Tier 1 pack (upload any or all)",
        "",
        f"- `{stem}.md` — long Markdown compendium",
        f"- `{stem}.html` — HTML export",
        f"- `{stem}.txt` — plain-text export",
        f"- `{stem}.docx` — Word",
        f"- `{stem}.pdf` — multi-page PDF compendium",
        "",
        "Other PDFs in this folder may come from `scripts/write_demo_pdf_assets_to_docs.py`; those are optional extras.",
        "",
        "## Tier 1 — retrieval / chat smoke tests",
        "",
    ]
    for i, q in enumerate(questions, start=1):
        lines.append(f"{i}. {q}")

    lines.extend(
        [
            "",
            "## Tier 2 pack (slides, spreadsheets, CSV, RTF)",
            "",
            f"- `{tier2_stem}.pptx` — PowerPoint deck",
            f"- `{tier2_stem}.xlsx` — Excel workbook (multiple sheets)",
            f"- `{tier2_stem}.csv` — tabular extract",
            f"- `{tier2_stem}.rtf` — rich text export",
        ]
    )
    if tier2_xls_written:
        lines.append(f"- `{tier2_stem}.xls` — legacy Excel (optional generator; requires `xlwt`)")
    else:
        lines.append("- *(No `.xls` file generated — `pip install xlwt` and re-run script to add it.)*")
    lines.extend(["", "## Tier 2 — retrieval / chat smoke tests", ""])
    for i, q in enumerate(tier2_q, start=1):
        lines.append(f"{i}. {q}")

    lines.extend(
        [
            "",
            "## Tier 3 pack (email, ebooks, OCR/image)",
            "",
            f"- `{tier3_stem}.eml` — email with text body + attachment marker content",
            f"- `{tier3_stem}.epub` — EPUB chapter sample",
            f"- `{tier3_stem}_ocr.png` — OCR image sample",
            f"- `{tier3_stem}_ocr.jpg` — OCR image sample (jpeg)",
            f"- `{tier3_stem}_ocr.tiff` — OCR image sample (tiff)",
            f"- `{tier3_stem}_scanned.pdf` — image-only PDF for OCR fallback testing",
        ]
    )
    if tier3_msg_written:
        lines.append(f"- `{tier3_stem}.msg` — Outlook .msg sample")
    else:
        lines.append("- *(No `.msg` generated on this machine. Create via Outlook Save As .msg or provide fixture.)*")
    if tier3_mobi_written:
        lines.append(f"- `{tier3_stem}.mobi` — MOBI sample")
    else:
        lines.append("- *(No `.mobi` generated on this machine. Install Calibre `ebook-convert` to auto-generate.)*")

    lines.extend(["", "## Tier 3 — retrieval / chat smoke tests", ""])
    for i, q in enumerate(tier3_q, start=1):
        lines.append(f"{i}. {q}")

    lines.extend(
        [
            "",
            "## Tips",
            "",
            "- Ask one question at a time; confirm citations point at the uploaded document.",
            "- If retrieval is empty, verify you selected the correct **workspace** and that ingestion finished.",
            "- Tier 2 uploads index text from slides, sheet rows, CSV lines, and RTF body — ask about strings you know appear in those files.",
            "- Tier 3 OCR quality depends on image clarity; for best results use high contrast black text on white background.",
            "",
        ]
    )
    return "\n".join(lines) + "\n"


def _paragraph(seed: int, topic: str) -> str:
    """Repeatable pseudo-prose (ASCII) for bulk size — not legal advice."""
    lines = [
        f"Sterling & Vale LLP internal reference — segment {seed} regarding {topic}. "
        "Partners and associates must align work product with the engagement letter, matter code, "
        "and workspace retention labels before circulation outside the firm.",
        f"Matter teams should calendar internal checkpoints at days {seed * 3}, {seed * 7}, and {seed * 11} "
        "from file opening for quality review, unless the client protocol specifies otherwise.",
        "Document naming: use SVL-[matter]-[doctype]-[version]; avoid client names in filenames on shared drives.",
        "Escalation: contact the practice group leader when exposure estimates exceed the matter budget band "
        "or when regulatory correspondence arrives from a new agency not listed in the intake form.",
    ]
    return lines[seed % len(lines)] + "\n\n"


def _large_markdown(workspace_name: str, tag: str, topic: str) -> str:
    header = f"# Sterling & Vale LLP — {workspace_name}\n\n"
    header += f"**Classification:** Internal use · **Workspace tag:** `{tag}` · **Focus:** {topic}\n\n"
    header += "---\n\n"
    body = ""
    for i in range(48):  # many sections -> larger files
        body += f"## {i + 1}. Practice note — {topic} (block {i + 1})\n\n"
        for j in range(4):
            body += _paragraph(i * 10 + j, topic)
        body += "### Checklist\n\n"
        body += (
            f"- [ ] Confirm matter code and workspace link for item {i + 1}\n"
            f"- [ ] Verify conflicts clearance timestamp SVL-{tag}-{i + 1:04d}\n"
            f"- [ ] Archive superseded drafts with `_deprecated` suffix\n\n"
        )
    body += "\n## Revision history\n\n| Version | Date | Owner |\n| --- | --- | --- |\n"
    for v in range(1, 9):
        body += f"| 2025.{v} | Q{v} | Knowledge steward |\n"
    return header + body


def _large_plaintext(workspace_name: str, tag: str, topic: str) -> str:
    text = (
        f"STERLING & VALE LLP — {workspace_name.upper()} — QUICK REFERENCE TEXT EXPORT\n"
        f"TAG={tag} TOPIC={topic}\n"
        + "=" * 72
        + "\n\n"
    )
    for i in range(60):
        text += f"[Section {i + 1}]\n"
        text += _paragraph(i, topic)
    return text


def _large_html(workspace_name: str, tag: str, topic: str) -> str:
    inner = ""
    for i in range(36):
        inner += f"<h2 id='s{i}'>Section {i + 1} — {html.escape(topic)}</h2>\n"
        inner += "<ul>\n"
        for j in range(3):
            inner += f"<li>{html.escape(_paragraph(i + j, topic).strip())}</li>\n"
        inner += "</ul>\n"
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8"/>
<title>Sterling &amp; Vale — {html.escape(workspace_name)}</title>
</head>
<body>
<header><h1>Sterling &amp; Vale LLP — {html.escape(workspace_name)}</h1>
<p><strong>Tag:</strong> {html.escape(tag)} · <strong>Topic:</strong> {html.escape(topic)}</p></header>
<main>
{inner}
</main>
</body>
</html>
"""


def _write_docx(path: Path, workspace_name: str, tag: str, topic: str) -> None:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(f"Sterling & Vale LLP — {workspace_name}", 0)
    p = doc.add_paragraph()
    p.add_run(f"Workspace tag {tag}. Focus: {topic}.").italic = True
    for i in range(40):
        doc.add_heading(f"Section {i + 1}: Operational detail", level=2)
        for j in range(3):
            doc.add_paragraph(_paragraph(i * 3 + j, topic).strip())
        doc.add_paragraph(
            f"Checklist SVL-{tag}-{i + 1:04d}: confirm intake form, matter code, and retention label.",
            style="List Bullet",
        )
    doc.save(path)


def _write_pdf(path: Path, title: str, workspace_name: str, tag: str, topic: str) -> None:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_left_margin(12)
    pdf.set_right_margin(12)
    pdf.add_page()
    w = pdf.w - pdf.l_margin - pdf.r_margin

    def emit_para(text: str, *, h: float, size: int = 9, style: str = "") -> None:
        pdf.set_font("Helvetica", style, size)
        safe = _fpdf_safe(text)
        for line in safe.splitlines():
            line = line.strip()
            if not line:
                pdf.ln(1)
                continue
            pdf.multi_cell(w, h, line)

    emit_para(title, h=7, size=13, style="B")
    pdf.ln(2)
    emit_para(f"Sterling & Vale LLP | {workspace_name} | tag {tag} | {topic}", h=4, size=9)
    pdf.ln(2)
    for i in range(240):  # multi-page compendium (larger on-disk size for upload tests)
        pdf.set_font("Helvetica", "B", 10)
        pdf.multi_cell(w, 5, _fpdf_safe(f"Section {i + 1}"))
        pdf.ln(1)
        pdf.set_font("Helvetica", size=9)
        emit_para(_paragraph(i, topic).strip(), h=4, size=9)
        pdf.ln(1)
    pdf.output(str(path))


def _plain_to_rtf(plain: str) -> str:
    esc = plain.replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}")
    return "{\\rtf1\\ansi\\deff0 {\\fonttbl{\\f0 Helvetica;}}\\f0\\fs22 " + esc.replace("\n", "\\par ") + "}"


def _write_tier2_pptx(path: Path, workspace_name: str, tag: str, topic: str) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"SVL Tier 2 - {workspace_name}"
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"Tag {tag} | {topic[:220]}"

    layout = prs.slide_layouts[1]
    for i in range(5):
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = f"Slide deck block {i + 1} - {tag}"
        tf = s.shapes.placeholders[1].text_frame
        tf.text = (
            _paragraph(i * 2, topic).strip()[:700]
            + f"\n\nInternal checkpoint reference: days {i * 3}, {i * 7}, {i * 11} from file opening."
        )
    prs.save(path)


def _write_tier2_xlsx(path: Path, workspace_name: str, tag: str, topic: str) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Summary"
    ws["A1"] = "Sterling Vale Tier 2 Summary"
    ws["A2"] = "Workspace"
    ws["B2"] = workspace_name
    ws["A3"] = f"SVL-TIER2-{tag}-XLSX-KEY"
    ws["B3"] = topic[:300]
    ws2 = wb.create_sheet("Checklists")
    ws2["A1"] = "Item"
    ws2["B1"] = "Value"
    ws2["A2"] = "Conflicts / clearance key"
    ws2["B2"] = f"SVL-TIER2-{tag}-KEY"
    ws2["A3"] = "Notes"
    ws2["B3"] = f"Verify conflicts clearance SVL-{tag}-T2 before substantive work; clause bank per workspace topic."
    ws3 = wb.create_sheet("Escalations")
    ws3["A1"] = "Escalation"
    ws3["A2"] = (
        "Contact the practice group leader when exposure exceeds the matter budget band, "
        "or for wall memo documentation within two business days where intake policy requires."
    )
    wb.save(path)


def _write_tier2_csv(path: Path, workspace_name: str, tag: str, topic: str) -> None:
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["record_id", "workspace_tag", "topic_line", "tier2_marker"])
        for i in range(45):
            w.writerow(
                [
                    f"SVL-{tag}-T2-{i + 1:04d}",
                    tag,
                    topic[:120],
                    f"checkpoint_band_{i % 5 + 1}",
                ]
            )


def _write_tier2_rtf(path: Path, workspace_name: str, tag: str, topic: str) -> None:
    plain = (
        f"Sterling Tier 2 RTF export for {workspace_name}. Workspace tag {tag}. "
        f"Focus: {topic}. "
    )
    plain += " ".join(_paragraph(k, topic).strip() for k in range(10))
    path.write_text(_plain_to_rtf(plain), encoding="utf-8")


def _write_tier2_xls(path: Path, workspace_name: str, tag: str, topic: str) -> bool:
    try:
        import xlwt
    except ImportError:
        return False
    wb = xlwt.Workbook()
    ws = wb.add_sheet("Tier2")
    ws.write(0, 0, "Sterling Vale XLS (legacy Tier 2)")
    ws.write(1, 0, "Workspace")
    ws.write(1, 1, workspace_name)
    ws.write(2, 0, f"SVL-TIER2-{tag}-XLS-KEY")
    ws.write(3, 0, topic[:250])
    for r in range(4, 30):
        ws.write(r, 0, _paragraph(r, topic)[:240])
    wb.save(path)
    return True


def _write_tier3_eml(path: Path, *, workspace_name: str, tag: str, topic: str, tier3_stem: str) -> None:
    msg = EmailMessage()
    msg["From"] = "knowledge-demo@sterlingvale.example"
    msg["To"] = "team@sterlingvale.example"
    msg["Subject"] = f"Tier 3 Demo Mail - {workspace_name}"
    msg.set_content(
        (
            f"Tier 3 EML body for {workspace_name}. Tag {tag}. Focus: {topic}.\n\n"
            "This message is used to verify header/body extraction and attachment parsing in ingestion.\n"
        )
    )
    attachment_text = (
        f"Attachment marker SVL-{tag}-T3-EMAIL-ATTACH.\n"
        f"Workspace={workspace_name}\n"
        "Use this marker in retrieval tests."
    )
    msg.add_attachment(
        attachment_text,
        subtype="plain",
        filename=f"{tier3_stem}_attachment.txt",
    )
    path.write_bytes(msg.as_bytes())


def _write_tier3_epub(path: Path, *, workspace_name: str, tag: str, topic: str) -> None:
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier(f"svl-tier3-{tag.lower()}")
    book.set_title(f"Sterling Vale Tier 3 EPUB - {workspace_name}")
    book.set_language("en")
    chapter = epub.EpubHtml(title="Tier 3 Intro", file_name="intro.xhtml", lang="en")
    chapter.content = (
        f"<h1>Tier 3 EPUB Intro - {workspace_name}</h1>"
        f"<p>Tag {tag}. Focus: {topic}.</p>"
        f"<p>Marker SVL-{tag}-T3-EPUB-101 for retrieval testing.</p>"
    )
    book.add_item(chapter)
    book.toc = (chapter,)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", chapter]
    epub.write_epub(str(path), book)


def _write_tier3_ocr_images(base: Path, *, workspace_name: str, tag: str, topic: str) -> list[Path]:
    from PIL import Image, ImageDraw

    text = (
        f"Sterling Vale Tier3 OCR\n"
        f"Workspace: {workspace_name}\n"
        f"Tag: {tag}\n"
        f"Marker: SVL-{tag}-T3-OCR-555\n"
        f"Topic: {topic[:90]}"
    )
    img = Image.new("RGB", (1200, 700), color="white")
    draw = ImageDraw.Draw(img)
    draw.multiline_text((40, 40), text, fill="black", spacing=12)

    outputs = [
        base.with_name(base.name + "_ocr.png"),
        base.with_name(base.name + "_ocr.jpg"),
        base.with_name(base.name + "_ocr.tiff"),
    ]
    img.save(outputs[0], format="PNG")
    img.save(outputs[1], format="JPEG", quality=94)
    img.save(outputs[2], format="TIFF")
    return outputs


def _write_tier3_scanned_pdf(path: Path, *, image_source: Path) -> None:
    from PIL import Image

    img = Image.open(image_source).convert("RGB")
    img.save(path, "PDF", resolution=200.0)


def _try_write_tier3_mobi(path: Path, *, source_epub: Path) -> bool:
    """
    Optional MOBI generation via Calibre's `ebook-convert` if installed.
    """
    ebook_convert = shutil.which("ebook-convert")
    if not ebook_convert:
        # Common Windows Calibre install paths.
        for candidate in (
            Path("C:/Program Files/Calibre2/ebook-convert.exe"),
            Path("C:/Program Files (x86)/Calibre2/ebook-convert.exe"),
            Path.home() / "AppData/Local/Programs/Calibre/ebook-convert.exe",
        ):
            if candidate.exists():
                ebook_convert = str(candidate)
                break
    if not ebook_convert:
        return False
    try:
        subprocess.run(
            [ebook_convert, str(source_epub), str(path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
    except Exception:
        return False
    return path.exists() and path.stat().st_size > 0


def _try_write_tier3_msg(path: Path, *, workspace_name: str, tag: str, topic: str) -> bool:
    """
    Optional Outlook COM generation on Windows hosts with Outlook installed.
    """
    ps = (
        "$ErrorActionPreference='Stop'; "
        "try { "
        "$outlook = New-Object -ComObject Outlook.Application; "
        "$mail = $outlook.CreateItem(0); "
        f"$mail.Subject = 'Tier 3 Demo MSG - {workspace_name}'; "
        f"$mail.Body = 'Workspace: {workspace_name}`nTag: {tag}`nMarker: SVL-{tag}-T3-MSG-777`nTopic: {topic}'; "
        f"$mail.SaveAs('{str(path).replace('\\', '\\\\')}', 3); "
        "$mail.Close(1); $outlook.Quit(); "
        "Write-Output 'OK'; "
        "} catch { Write-Output 'NO'; }"
    )
    try:
        r = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
    except Exception:
        return False
    return ("OK" in (r.stdout or "")) and path.exists() and path.stat().st_size > 0


def main() -> None:
    base = _REPO / "docs" / "demo" / "organizations" / DEMO_ORGANIZATION_SLUG / "workspaces"
    base.mkdir(parents=True, exist_ok=True)

    for workspace_name, tag, topic in _WORKSPACES:
        slug = workspace_dir_slug(workspace_name)
        dest_dir = base / slug
        dest_dir.mkdir(parents=True, exist_ok=True)

        stem = f"Sterling_Vale_{tag}_Demo_Pack"
        tier2_stem = f"Sterling_Vale_{tag}_Tier2_Pack"
        tier3_stem = f"Sterling_Vale_{tag}_Tier3_Pack"

        md_path = dest_dir / f"{stem}.md"
        if not md_path.exists():
            md_path.write_text(_large_markdown(workspace_name, tag, topic), encoding="utf-8")

        txt_path = dest_dir / f"{stem}.txt"
        if not txt_path.exists():
            txt_path.write_text(_large_plaintext(workspace_name, tag, topic), encoding="utf-8")

        html_path = dest_dir / f"{stem}.html"
        if not html_path.exists():
            html_path.write_text(_large_html(workspace_name, tag, topic), encoding="utf-8")

        docx_path = dest_dir / f"{stem}.docx"
        if not docx_path.exists():
            _write_docx(docx_path, workspace_name, tag, topic)

        pdf_title = f"Sterling & Vale - {workspace_name} - Demo Compendium"
        pdf_path = dest_dir / f"{stem}.pdf"
        if not pdf_path.exists():
            _write_pdf(pdf_path, pdf_title, workspace_name, tag, topic)

        tier2_pptx = dest_dir / f"{tier2_stem}.pptx"
        if not tier2_pptx.exists():
            _write_tier2_pptx(tier2_pptx, workspace_name, tag, topic)
        tier2_xlsx = dest_dir / f"{tier2_stem}.xlsx"
        if not tier2_xlsx.exists():
            _write_tier2_xlsx(tier2_xlsx, workspace_name, tag, topic)
        tier2_csv = dest_dir / f"{tier2_stem}.csv"
        if not tier2_csv.exists():
            _write_tier2_csv(tier2_csv, workspace_name, tag, topic)
        tier2_rtf = dest_dir / f"{tier2_stem}.rtf"
        if not tier2_rtf.exists():
            _write_tier2_rtf(tier2_rtf, workspace_name, tag, topic)
        tier2_xls_path = dest_dir / f"{tier2_stem}.xls"
        tier2_xls_written = tier2_xls_path.exists() or _write_tier2_xls(tier2_xls_path, workspace_name, tag, topic)

        # Tier 3 sample pack.
        tier3_eml = dest_dir / f"{tier3_stem}.eml"
        _write_tier3_eml(tier3_eml, workspace_name=workspace_name, tag=tag, topic=topic, tier3_stem=tier3_stem)
        tier3_epub = dest_dir / f"{tier3_stem}.epub"
        _write_tier3_epub(tier3_epub, workspace_name=workspace_name, tag=tag, topic=topic)
        tier3_img_base = dest_dir / tier3_stem
        tier3_imgs = _write_tier3_ocr_images(tier3_img_base, workspace_name=workspace_name, tag=tag, topic=topic)
        _write_tier3_scanned_pdf(dest_dir / f"{tier3_stem}_scanned.pdf", image_source=tier3_imgs[0])
        tier3_mobi_path = dest_dir / f"{tier3_stem}.mobi"
        tier3_mobi_written = tier3_mobi_path.exists() or _try_write_tier3_mobi(
            tier3_mobi_path, source_epub=tier3_epub
        )
        tier3_msg_path = dest_dir / f"{tier3_stem}.msg"
        tier3_msg_written = tier3_msg_path.exists() or _try_write_tier3_msg(
            tier3_msg_path,
            workspace_name=workspace_name,
            tag=tag,
            topic=topic,
        )

        sq_path = dest_dir / "Sample_Questions.md"
        sq_path.write_text(
            _sample_questions_markdown(
                workspace_name,
                tag,
                topic,
                stem,
                tier2_stem,
                tier3_stem,
                tier2_xls_written=tier2_xls_written,
                tier3_msg_written=tier3_msg_written,
                tier3_mobi_written=tier3_mobi_written,
            ),
            encoding="utf-8",
        )

        n_pack = len(list(dest_dir.glob(f"{stem}.*")))
        n_t2 = len(list(dest_dir.glob(f"{tier2_stem}.*")))
        n_t3 = len(list(dest_dir.glob(f"{tier3_stem}*")))
        print(
            f"Wrote Tier1 pack ({n_pack}) + Tier2 pack ({n_t2}) + Tier3 pack ({n_t3}) + Sample_Questions.md "
            f"-> {dest_dir.relative_to(_REPO)}"
        )

    print("Done.")


if __name__ == "__main__":
    main()
