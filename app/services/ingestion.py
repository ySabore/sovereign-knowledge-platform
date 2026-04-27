from __future__ import annotations

import csv
import email
import hashlib
import io
import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from uuid import UUID, uuid4

from fastapi import HTTPException, UploadFile, status
from pypdf import PdfReader

from app.config import settings
from app.services.storage import get_storage_backend
from app.services.text_cleaner import clean_text_for_ingestion

# Tier 1 RAG uploads: PDF, Word, text, Markdown, HTML (see docs / product parity).
TIER1_UPLOAD_EXTENSIONS: frozenset[str] = frozenset({".pdf", ".docx", ".txt", ".md", ".markdown", ".html", ".htm"})
# Tier 2: slides, spreadsheets, CSV, RTF.
TIER2_UPLOAD_EXTENSIONS: frozenset[str] = frozenset({".pptx", ".xlsx", ".xls", ".csv", ".rtf"})
# Tier 3: email + images + OCR fallback for scanned PDFs (when dependencies exist).
TIER3_UPLOAD_EXTENSIONS: frozenset[str] = frozenset(
    {".eml", ".msg", ".epub", ".mobi", ".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}
)
ALLOWED_UPLOAD_EXTENSIONS: frozenset[str] = TIER1_UPLOAD_EXTENSIONS | TIER2_UPLOAD_EXTENSIONS | TIER3_UPLOAD_EXTENSIONS
_WHITESPACE_RE = re.compile(r"\s+")
_MD_HEADER = re.compile(r"^\s{0,3}(#{1,6})\s+(.+?)\s*$")
_PARA_SPLIT = re.compile(r"\n\s*\n+")
# Attachment extraction guardrails: avoid unbounded recursion or huge nested payloads.
EMAIL_ATTACHMENT_MAX_BYTES = 5 * 1024 * 1024
EMAIL_ATTACHMENT_MAX_DEPTH = 2


@dataclass(slots=True)
class ExtractedPage:
    page_number: int
    text: str


@dataclass(slots=True)
class ChunkRecord:
    chunk_index: int
    page_number: int | None
    section_title: str | None
    content: str
    char_count: int


@dataclass(slots=True)
class StoredUpload:
    storage_path: str
    extraction_path: str
    checksum_sha256: str
    size_bytes: int
    storage_provider: str | None = None
    storage_bucket: str | None = None
    storage_key: str | None = None
    storage_etag: str | None = None


def validate_document_upload(upload: UploadFile) -> None:
    filename = (upload.filename or "").strip()
    if not filename:
        raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail="Filename is required")
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_UPLOAD_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
            detail=(
                "Unsupported file type. Supported: Tier 1 — PDF (.pdf), Word (.docx), plain text (.txt), "
                "Markdown (.md, .markdown), HTML (.html, .htm); Tier 2 — PowerPoint (.pptx), Excel (.xlsx, .xls), "
                "CSV (.csv), RTF (.rtf); Tier 3 — email (.eml, .msg), eBooks (.epub, .mobi), "
                "images (.png, .jpg, .jpeg, .webp, .tif, .tiff). "
                "OCR is used for images and scanned PDFs when available."
            ),
        )


async def persist_upload_file(upload: UploadFile, storage_root: Path, workspace_id: UUID) -> StoredUpload:
    validate_document_upload(upload)
    safe_name = Path(upload.filename or "document.pdf").name
    tmp_dir = storage_root / "_tmp_uploads"
    tmp_dir.mkdir(parents=True, exist_ok=True)
    destination_path = tmp_dir / f"{uuid4()}-{safe_name}"

    hasher = hashlib.sha256()
    size_bytes = 0
    max_bytes = settings.max_upload_size_bytes
    with destination_path.open("wb") as out:
        while True:
            chunk = await upload.read(1024 * 1024)
            if not chunk:
                break
            size_bytes += len(chunk)
            if size_bytes > max_bytes:
                out.close()
                destination_path.unlink(missing_ok=True)
                await upload.close()
                raise HTTPException(
                    status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                    detail=f"Upload exceeds maximum size of {settings.max_upload_size_mb} MB",
                )
            out.write(chunk)
            hasher.update(chunk)

    await upload.close()
    backend = get_storage_backend()
    stored = backend.store_upload(
        local_path=destination_path,
        workspace_id=workspace_id,
        safe_name=safe_name,
        checksum_sha256=hasher.hexdigest(),
        size_bytes=size_bytes,
    )
    return StoredUpload(
        storage_path=stored.storage_uri,
        extraction_path=stored.extraction_path,
        checksum_sha256=stored.checksum_sha256,
        size_bytes=stored.size_bytes,
        storage_provider=stored.provider,
        storage_bucket=stored.bucket,
        storage_key=stored.key,
        storage_etag=stored.etag,
    )


def extract_pdf_pages(file_path: str) -> list[ExtractedPage]:
    reader = PdfReader(file_path)
    pages: list[ExtractedPage] = []
    for index, page in enumerate(reader.pages, start=1):
        raw_text = page.extract_text() or ""
        normalized = normalize_text(raw_text)
        cleaned = clean_text_for_ingestion(normalized, strip_html=False)
        if cleaned:
            pages.append(ExtractedPage(page_number=index, text=cleaned))
    return pages


def _ocr_available() -> bool:
    try:
        import pytesseract  # noqa: F401

        return True
    except Exception:
        return False


def _extract_image_ocr_pages(file_path: str) -> list[ExtractedPage]:
    """
    OCR image formats (png/jpg/webp/tiff). Requires `pytesseract` and a `tesseract` binary.
    """
    try:
        import pytesseract
        from PIL import Image
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Image OCR requires pytesseract + Pillow on the API server.",
        ) from exc

    try:
        img = Image.open(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable image: {exc}",
        ) from exc

    try:
        text = pytesseract.image_to_string(img)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=f"OCR failed (tesseract unavailable or errored): {exc}",
        ) from exc

    cleaned = clean_text_for_ingestion(text, strip_html=False)
    if not cleaned.strip():
        return []
    return [ExtractedPage(page_number=1, text=cleaned)]


def _extract_pdf_ocr_pages(file_path: str, *, max_pages: int = 40, dpi: int = 200) -> list[ExtractedPage]:
    """
    OCR image-only PDFs. Requires `pymupdf` (fitz) to render and pytesseract + tesseract to OCR.
    """
    try:
        import fitz  # PyMuPDF
        import pytesseract
        from PIL import Image
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="PDF OCR requires pymupdf + pytesseract + Pillow on the API server.",
        ) from exc

    try:
        doc = fitz.open(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable PDF: {exc}",
        ) from exc

    pages: list[ExtractedPage] = []
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    try:
        for idx in range(min(len(doc), max_pages)):
            page = doc.load_page(idx)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            try:
                text = pytesseract.image_to_string(img)
            except Exception as exc:
                raise HTTPException(
                    status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                    detail=f"OCR failed (tesseract unavailable or errored): {exc}",
                ) from exc
            cleaned = clean_text_for_ingestion(text, strip_html=False)
            if cleaned.strip():
                pages.append(ExtractedPage(page_number=idx + 1, text=cleaned))
    finally:
        doc.close()
    return pages


def _extract_eml_pages(file_path: str, *, depth: int = 0) -> list[ExtractedPage]:
    """
    Parse RFC822 `.eml` email. Prefer text/plain, fall back to text/html.
    """
    raw = Path(file_path).read_bytes()
    try:
        msg = email.message_from_bytes(raw)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable email (.eml): {exc}",
        ) from exc

    subject = (msg.get("subject") or "").strip()
    from_ = (msg.get("from") or "").strip()
    to_ = (msg.get("to") or "").strip()
    date_ = (msg.get("date") or "").strip()

    plain_parts: list[str] = []
    html_parts: list[str] = []
    attachment_sections: list[str] = []

    if msg.is_multipart():
        for part in msg.walk():
            ctype = (part.get_content_type() or "").lower()
            disp = (part.get("content-disposition") or "").lower()
            filename = (part.get_filename() or "").strip()
            if "attachment" in disp or filename:
                payload = part.get_payload(decode=True) or b""
                att_text = _extract_attachment_text_from_bytes(
                    filename=filename or "attachment",
                    payload=payload,
                    content_type=ctype,
                    depth=depth,
                )
                if att_text:
                    attachment_sections.append(att_text)
                continue
            if ctype not in ("text/plain", "text/html"):
                continue
            payload = part.get_payload(decode=True) or b""
            charset = part.get_content_charset() or "utf-8"
            try:
                text = payload.decode(charset, errors="replace")
            except Exception:
                text = payload.decode("utf-8", errors="replace")
            if ctype == "text/plain":
                plain_parts.append(text)
            else:
                html_parts.append(text)
    else:
        payload = msg.get_payload(decode=True) or b""
        charset = msg.get_content_charset() or "utf-8"
        try:
            text = payload.decode(charset, errors="replace")
        except Exception:
            text = payload.decode("utf-8", errors="replace")
        ctype = (msg.get_content_type() or "").lower()
        if ctype == "text/html":
            html_parts.append(text)
        else:
            plain_parts.append(text)

    body_plain = "\n\n".join(t.strip() for t in plain_parts if (t or "").strip())
    body_html = "\n\n".join(t.strip() for t in html_parts if (t or "").strip())

    header = "\n".join(
        s
        for s in [
            f"Subject: {subject}" if subject else "",
            f"From: {from_}" if from_ else "",
            f"To: {to_}" if to_ else "",
            f"Date: {date_}" if date_ else "",
        ]
        if s
    )
    assembled = (header + "\n\n" + (body_plain or "")).strip()
    if not assembled and body_html:
        assembled = (header + "\n\n" + body_html).strip()
        cleaned = clean_text_for_ingestion(assembled, strip_html=True)
    else:
        cleaned = clean_text_for_ingestion(assembled, strip_html=False)
        if not cleaned.strip() and body_html:
            cleaned = clean_text_for_ingestion((header + "\n\n" + body_html).strip(), strip_html=True)

    if attachment_sections:
        with_atts = (cleaned + "\n\n" + "\n\n".join(attachment_sections)).strip() if cleaned else "\n\n".join(attachment_sections)
        cleaned = clean_text_for_ingestion(with_atts, strip_html=False)

    if not cleaned.strip():
        return []
    return [ExtractedPage(page_number=1, text=cleaned)]


def _extract_attachment_text_from_bytes(*, filename: str, payload: bytes, content_type: str, depth: int) -> str | None:
    """
    Best-effort extraction for email attachments (text-first, bounded size).
    Also performs bounded recursive extraction for supported binary formats.
    """
    name = (filename or "attachment").strip()
    ext = Path(name).suffix.lower()
    if not payload or len(payload) > EMAIL_ATTACHMENT_MAX_BYTES:
        return None

    prefix = f"Attachment: {name}\n"
    if ext in (".txt", ".md", ".markdown"):
        txt = payload.decode("utf-8", errors="replace")
        body = clean_text_for_ingestion(txt, strip_html=False)
        return (prefix + body).strip() if body.strip() else None
    if ext in (".html", ".htm"):
        txt = payload.decode("utf-8", errors="replace")
        body = clean_text_for_ingestion(txt, strip_html=True)
        return (prefix + body).strip() if body.strip() else None
    if ext == ".csv":
        decoded = _decode_raw_for_csv(payload)
        rows: list[str] = []
        for row in csv.reader(io.StringIO(decoded)):
            vals = [(c or "").strip() for c in row]
            line = " | ".join(c for c in vals if c)
            if line:
                rows.append(line)
        body = clean_text_for_ingestion("\n".join(rows), strip_html=False)
        return (prefix + body).strip() if body.strip() else None
    if ext == ".rtf":
        try:
            from striprtf.striprtf import rtf_to_text
        except Exception:
            return None
        txt = payload.decode("utf-8", errors="replace")
        try:
            plain = rtf_to_text(txt)
        except Exception:
            return None
        body = clean_text_for_ingestion(plain, strip_html=False)
        return (prefix + body).strip() if body.strip() else None
    if ext == ".eml":
        try:
            nested = email.message_from_bytes(payload)
        except Exception:
            return None
        subj = (nested.get("subject") or "").strip()
        from_ = (nested.get("from") or "").strip()
        to_ = (nested.get("to") or "").strip()
        text_parts: list[str] = []
        if nested.is_multipart():
            for p in nested.walk():
                if (p.get_content_type() or "").lower() != "text/plain":
                    continue
                pdata = p.get_payload(decode=True) or b""
                text_parts.append(pdata.decode(p.get_content_charset() or "utf-8", errors="replace"))
        else:
            pdata = nested.get_payload(decode=True) or b""
            text_parts.append(pdata.decode(nested.get_content_charset() or "utf-8", errors="replace"))
        body = clean_text_for_ingestion(
            "\n".join(s for s in [f"Subject: {subj}", f"From: {from_}", f"To: {to_}", "\n".join(text_parts)] if s),
            strip_html=False,
        )
        return (prefix + body).strip() if body.strip() else None
    # Recursive extraction for binary attachments already supported by main ingestion path.
    if ext in ALLOWED_UPLOAD_EXTENSIONS and depth < EMAIL_ATTACHMENT_MAX_DEPTH:
        nested_pages: list[ExtractedPage] = []
        try:
            with tempfile.TemporaryDirectory(prefix="att-") as tmp:
                p = Path(tmp) / name
                p.write_bytes(payload)
                nested_pages = _extract_pages_from_upload(str(p), name, depth=depth + 1)
        except Exception:
            nested_pages = []
        if nested_pages:
            nested_text = "\n\n".join(pg.text for pg in nested_pages if (pg.text or "").strip())
            body = clean_text_for_ingestion(nested_text, strip_html=False)
            return (prefix + body).strip() if body.strip() else None
    # Skip binary Office/PDF/image attachments in this step (separate doc ingest path).
    if content_type.startswith("text/"):
        txt = payload.decode("utf-8", errors="replace")
        body = clean_text_for_ingestion(txt, strip_html=False)
        return (prefix + body).strip() if body.strip() else None
    return None


def _extract_msg_pages(file_path: str, *, depth: int = 0) -> list[ExtractedPage]:
    """
    Parse Outlook `.msg` email via extract_msg.
    """
    try:
        import extract_msg
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Outlook .msg ingestion requires extract-msg on the API server.",
        ) from exc

    try:
        msg = extract_msg.Message(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable Outlook email (.msg): {exc}",
        ) from exc

    subject = (msg.subject or "").strip()
    sender = (msg.sender or "").strip()
    to_ = (msg.to or "").strip()
    date_ = (str(msg.date) if msg.date is not None else "").strip()
    body_plain = (msg.body or "").strip()
    body_html = getattr(msg, "htmlBody", None)
    if isinstance(body_html, bytes):
        body_html = body_html.decode("utf-8", errors="replace")
    if not isinstance(body_html, str):
        body_html = ""

    header = "\n".join(
        s
        for s in [
            f"Subject: {subject}" if subject else "",
            f"From: {sender}" if sender else "",
            f"To: {to_}" if to_ else "",
            f"Date: {date_}" if date_ else "",
        ]
        if s
    )

    assembled = (header + "\n\n" + body_plain).strip()
    cleaned = clean_text_for_ingestion(assembled, strip_html=False)
    if not cleaned.strip() and body_html:
        cleaned = clean_text_for_ingestion((header + "\n\n" + body_html).strip(), strip_html=True)

    attachment_sections: list[str] = []
    for att in getattr(msg, "attachments", []) or []:
        name = (getattr(att, "longFilename", None) or getattr(att, "filename", None) or "attachment").strip()
        data = getattr(att, "data", None)
        if isinstance(data, str):
            payload = data.encode("utf-8", errors="replace")
        elif isinstance(data, bytes):
            payload = data
        else:
            payload = b""
        ctype = (getattr(att, "mimetype", None) or "").lower()
        att_text = _extract_attachment_text_from_bytes(filename=name, payload=payload, content_type=ctype, depth=depth)
        if att_text:
            attachment_sections.append(att_text)

    if attachment_sections:
        with_atts = (cleaned + "\n\n" + "\n\n".join(attachment_sections)).strip() if cleaned else "\n\n".join(attachment_sections)
        cleaned = clean_text_for_ingestion(with_atts, strip_html=False)

    try:
        msg.close()
    except Exception:
        pass

    if not cleaned.strip():
        return []
    return [ExtractedPage(page_number=1, text=cleaned)]


def _extract_epub_pages(file_path: str) -> list[ExtractedPage]:
    """
    Parse EPUB chapters/documents and return one page-like segment per document item.
    """
    try:
        from bs4 import BeautifulSoup
        from ebooklib import ITEM_DOCUMENT, epub
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="EPUB ingestion requires ebooklib + beautifulsoup4 on the API server.",
        ) from exc

    try:
        book = epub.read_epub(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable EPUB: {exc}",
        ) from exc

    pages: list[ExtractedPage] = []
    page_no = 1
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        payload = item.get_body_content() or item.get_content() or b""
        html_text = payload.decode("utf-8", errors="replace")
        text = BeautifulSoup(html_text, "html.parser").get_text("\n")
        cleaned = clean_text_for_ingestion(text, strip_html=False)
        if cleaned.strip():
            pages.append(ExtractedPage(page_number=page_no, text=cleaned))
            page_no += 1
    return pages


def _extract_mobi_pages(file_path: str) -> list[ExtractedPage]:
    """
    Best-effort MOBI extraction using `mobi` package to unpack, then parse HTML/XHTML text.
    """
    try:
        from bs4 import BeautifulSoup
        import mobi
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="MOBI ingestion requires mobi + beautifulsoup4 on the API server.",
        ) from exc

    extracted_tmp: str | None = None
    try:
        extracted_tmp, _ = mobi.extract(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable MOBI: {exc}",
        ) from exc

    pages: list[ExtractedPage] = []
    page_no = 1
    try:
        root = Path(extracted_tmp or tempfile.mkdtemp(prefix="mobi-empty-"))
        for p in sorted(root.rglob("*")):
            if p.suffix.lower() not in (".html", ".htm", ".xhtml", ".xml"):
                continue
            html_text = p.read_text(encoding="utf-8", errors="replace")
            text = BeautifulSoup(html_text, "html.parser").get_text("\n")
            cleaned = clean_text_for_ingestion(text, strip_html=False)
            if cleaned.strip():
                pages.append(ExtractedPage(page_number=page_no, text=cleaned))
                page_no += 1
    finally:
        try:
            if extracted_tmp:
                mobi.cleanup(extracted_tmp)
        except Exception:
            pass
    return pages


def _extract_docx_pages(file_path: str) -> list[ExtractedPage]:
    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable Word document (.docx): {exc}",
        ) from exc
    parts: list[str] = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join((c.text or "").strip() for c in row.cells)
            if row_text.strip():
                parts.append(row_text)
    body = "\n\n".join(parts)
    cleaned = clean_text_for_ingestion(body, strip_html=False)
    if not cleaned.strip():
        return []
    return [ExtractedPage(page_number=1, text=cleaned)]


def _extract_utf8_text_pages(file_path: str, *, strip_html: bool) -> list[ExtractedPage]:
    raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
    cleaned = clean_text_for_ingestion(raw, strip_html=strip_html)
    if not cleaned.strip():
        return []
    return [ExtractedPage(page_number=1, text=cleaned)]


def _shape_text_pptx(shape: object) -> list[str]:
    """Collect visible text from a slide shape (groups, text frames, tables)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out: list[str] = []
    st = getattr(shape, "shape_type", None)
    if st == MSO_SHAPE_TYPE.GROUP:
        for child in getattr(shape, "shapes", []):
            out.extend(_shape_text_pptx(child))
        return out
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                out.append(line)
        return out
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            line = "".join((run.text or "") for run in para.runs).strip()
            if line:
                out.append(line)
    return out


def _extract_pptx_pages(file_path: str) -> list[ExtractedPage]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="PowerPoint ingestion requires python-pptx on the API server.",
        ) from exc
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable PowerPoint (.pptx): {exc}",
        ) from exc
    pages: list[ExtractedPage] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_text_pptx(shape))
        body = "\n\n".join(parts)
        cleaned = clean_text_for_ingestion(body, strip_html=False)
        if cleaned.strip():
            pages.append(ExtractedPage(page_number=slide_idx, text=cleaned))
    return pages


def _extract_xlsx_pages(file_path: str) -> list[ExtractedPage]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Excel .xlsx ingestion requires openpyxl on the API server.",
        ) from exc
    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable Excel workbook (.xlsx): {exc}",
        ) from exc
    pages: list[ExtractedPage] = []
    try:
        for si, sheet in enumerate(wb, start=1):
            rows_out: list[str] = []
            rows_out.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() if c is not None else "" for c in row]
                line = " | ".join(c for c in cells if c)
                if line:
                    rows_out.append(line)
            body = "\n".join(rows_out)
            cleaned = clean_text_for_ingestion(body, strip_html=False)
            if cleaned.strip():
                pages.append(ExtractedPage(page_number=si, text=cleaned))
    finally:
        wb.close()
    return pages


def _extract_xls_pages(file_path: str) -> list[ExtractedPage]:
    try:
        import xlrd
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Excel .xls ingestion requires xlrd on the API server.",
        ) from exc
    try:
        book = xlrd.open_workbook(file_path)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable Excel workbook (.xls): {exc}",
        ) from exc
    pages: list[ExtractedPage] = []
    for si in range(book.nsheets):
        sheet = book.sheet_by_index(si)
        rows_out: list[str] = [f"Sheet: {sheet.name}"]
        for ri in range(sheet.nrows):
            vals = sheet.row_values(ri)
            cells = [str(c).strip() for c in vals if c is not None and str(c).strip()]
            if cells:
                rows_out.append(" | ".join(cells))
        body = "\n".join(rows_out)
        cleaned = clean_text_for_ingestion(body, strip_html=False)
        if cleaned.strip():
            pages.append(ExtractedPage(page_number=si + 1, text=cleaned))
    return pages


def _decode_raw_for_csv(raw: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_csv_pages(file_path: str) -> list[ExtractedPage]:
    raw = Path(file_path).read_bytes()
    text = _decode_raw_for_csv(raw)
    reader = csv.reader(io.StringIO(text))
    rows_out: list[str] = []
    for row in reader:
        cells = [(c or "").strip() for c in row]
        line = " | ".join(c for c in cells if c)
        if line:
            rows_out.append(line)
    body = "\n".join(rows_out)
    cleaned = clean_text_for_ingestion(body, strip_html=False)
    if not cleaned.strip():
        return []
    return [ExtractedPage(page_number=1, text=cleaned)]


def _extract_rtf_pages(file_path: str) -> list[ExtractedPage]:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="RTF ingestion requires striprtf on the API server.",
        ) from exc
    raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
    try:
        plain = rtf_to_text(raw)
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Invalid or unreadable RTF: {exc}",
        ) from exc
    cleaned = clean_text_for_ingestion(plain, strip_html=False)
    if not cleaned.strip():
        return []
    return [ExtractedPage(page_number=1, text=cleaned)]


def _extract_pages_from_upload(file_path: str, original_filename: str, *, depth: int) -> list[ExtractedPage]:
    """
    Dispatch by filename extension to produce page-like segments for chunking.
    PDFs and PPTX use one segment per page/slide; multi-sheet spreadsheets use one per sheet.
    """
    ext = Path(original_filename).suffix.lower()
    if ext not in ALLOWED_UPLOAD_EXTENSIONS:
        ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        pages = extract_pdf_pages(file_path)
        if pages:
            return pages
        # OCR fallback for scanned/image-only PDFs (best-effort when deps exist).
        if _ocr_available():
            try:
                return _extract_pdf_ocr_pages(file_path)
            except HTTPException:
                return []
        return []
    if ext == ".docx":
        return _extract_docx_pages(file_path)
    if ext in (".txt", ".md", ".markdown"):
        return _extract_utf8_text_pages(file_path, strip_html=False)
    if ext in (".html", ".htm"):
        return _extract_utf8_text_pages(file_path, strip_html=True)
    if ext == ".pptx":
        return _extract_pptx_pages(file_path)
    if ext == ".xlsx":
        return _extract_xlsx_pages(file_path)
    if ext == ".xls":
        return _extract_xls_pages(file_path)
    if ext == ".csv":
        return _extract_csv_pages(file_path)
    if ext == ".rtf":
        return _extract_rtf_pages(file_path)
    if ext == ".eml":
        return _extract_eml_pages(file_path, depth=depth)
    if ext == ".msg":
        return _extract_msg_pages(file_path, depth=depth)
    if ext == ".epub":
        return _extract_epub_pages(file_path)
    if ext == ".mobi":
        return _extract_mobi_pages(file_path)
    if ext in (".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"):
        return _extract_image_ocr_pages(file_path)
    return []


def extract_pages_from_upload(file_path: str, original_filename: str) -> list[ExtractedPage]:
    """
    Public extraction entrypoint.
    Uses internal bounded recursion for nested email attachments.
    """
    return _extract_pages_from_upload(file_path, original_filename, depth=0)


def normalize_text(value: str) -> str:
    return _WHITESPACE_RE.sub(" ", value).strip()


def _paragraphs(text: str) -> list[str]:
    parts = _PARA_SPLIT.split(text)
    return [p.strip() for p in parts if p.strip()]


def _sliding_chunks(text: str, *, chunk_size: int, step: int) -> list[str]:
    if not text:
        return []
    out: list[str] = []
    pos = 0
    while pos < len(text):
        piece = text[pos : pos + chunk_size].strip()
        if piece:
            out.append(piece)
        pos += step
    return out


def _effective_chunk_params() -> tuple[int, int]:
    """~400–600 token targets (chars/4 heuristic) with ~10% overlap by default."""
    chunk_size = max(400, int(settings.ingestion_target_tokens * 4))
    overlap = int(chunk_size * settings.ingestion_overlap_ratio)
    overlap = max(40, min(overlap, chunk_size - 1))
    return chunk_size, overlap


def build_chunks(
    pages: list[ExtractedPage],
    *,
    chunk_size: int | None = None,
    chunk_overlap: int | None = None,
) -> list[ChunkRecord]:
    if chunk_size is None or chunk_overlap is None:
        chunk_size, chunk_overlap = _effective_chunk_params()
    if chunk_size <= 0:
        raise ValueError("chunk_size must be > 0")
    if chunk_overlap < 0 or chunk_overlap >= chunk_size:
        raise ValueError("chunk_overlap must be >= 0 and < chunk_size")

    step = chunk_size - chunk_overlap
    chunks: list[ChunkRecord] = []
    chunk_index = 0

    for page in pages:
        current_section: str | None = None
        paras = _paragraphs(page.text)
        buf: list[str] = []
        buf_len = 0

        def flush_buf() -> None:
            nonlocal chunk_index, buf, buf_len
            if not buf:
                return
            body = "\n\n".join(buf)
            prefix = f"Section: {current_section}\n\n" if current_section else ""
            full = (prefix + body).strip()
            chunks.append(
                ChunkRecord(
                    chunk_index=chunk_index,
                    page_number=page.page_number,
                    section_title=current_section,
                    content=full,
                    char_count=len(full),
                )
            )
            chunk_index += 1
            buf = []
            buf_len = 0

        for para in paras:
            hm = _MD_HEADER.match(para)
            if hm:
                flush_buf()
                current_section = hm.group(2).strip()
                continue

            if len(para) > chunk_size:
                flush_buf()
                for piece in _sliding_chunks(para, chunk_size=chunk_size, step=step):
                    prefix = f"Section: {current_section}\n\n" if current_section else ""
                    full = (prefix + piece).strip()
                    chunks.append(
                        ChunkRecord(
                            chunk_index=chunk_index,
                            page_number=page.page_number,
                            section_title=current_section,
                            content=full,
                            char_count=len(full),
                        )
                    )
                    chunk_index += 1
                continue

            add_len = len(para) if not buf else len(para) + 2
            if buf_len + add_len <= chunk_size:
                buf.append(para)
                buf_len += add_len
            else:
                flush_buf()
                if len(para) > chunk_size:
                    for piece in _sliding_chunks(para, chunk_size=chunk_size, step=step):
                        prefix = f"Section: {current_section}\n\n" if current_section else ""
                        full = (prefix + piece).strip()
                        chunks.append(
                            ChunkRecord(
                                chunk_index=chunk_index,
                                page_number=page.page_number,
                                section_title=current_section,
                                content=full,
                                char_count=len(full),
                            )
                        )
                        chunk_index += 1
                else:
                    buf = [para]
                    buf_len = len(para)

        flush_buf()

    return chunks


def build_chunks_from_plain_text(
    text: str,
    *,
    chunk_size: int | None = None,
    chunk_overlap: int | None = None,
    page_number: int | None = None,
) -> list[ChunkRecord]:
    """Chunk arbitrary cleaned text (connector/HTML body) with the same paragraph + section strategy."""
    cleaned = clean_text_for_ingestion(text, strip_html=True)
    if not cleaned:
        return []
    page = ExtractedPage(page_number=page_number or 1, text=cleaned)
    return build_chunks([page], chunk_size=chunk_size, chunk_overlap=chunk_overlap)
